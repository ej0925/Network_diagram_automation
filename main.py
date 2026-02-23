from fastapi import FastAPI, Response
from pydantic import BaseModel
from typing import List, Optional
import pandas as pd
import networkx as nx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from collections import defaultdict
import io

app = FastAPI()

# 定義從 Power Automate 接收的 JSON 資料格式
class CircuitData(BaseModel):
    PRODUCT_TYPE: str
    CIRCUIT_ID: str
    SITE_ID_A: str
    BUILDING_A: str
    SITE_ID_B: Optional[str] = ""
    BUILDING_B: Optional[str] = ""
    ROUTER_NAME: Optional[str] = ""

# --- L2R 排版引擎 (定義在外面) ---
def apply_l2r_topology_layout(G):
    left_nodes = []   # 客戶地址
    center_nodes = [] # 網路/雲端
    right_nodes = []  # 機房/總部

    for n, d in G.nodes(data=True):
        role = d.get('role', 'customer')
        if role == 'cloud': center_nodes.append(n)
        elif role == 'datacenter': right_nodes.append(n)
        else: left_nodes.append(n)

    pos = {}
    X_LEFT, X_CENTER, X_RIGHT = 2.0, 6.66, 11.33

    def assign_y_positions(nodes, x_pos):
        if not nodes: return
        total_nodes = len(nodes)
        y_spacing = 1.5
        start_y = 3.75 - ((total_nodes - 1) * y_spacing / 2)
        for i, node in enumerate(nodes):
            pos[node] = [x_pos, start_y + (i * y_spacing)]

    assign_y_positions(left_nodes, X_LEFT)
    assign_y_positions(center_nodes, X_CENTER)
    assign_y_positions(right_nodes, X_RIGHT)
    return pos

# --- API 執行主體 ---
@app.post("/generate-pptx")
def generate_pptx(data: List[CircuitData]):
    # 1. 將收到的 JSON 轉成 DataFrame 
    # (註：若您使用最新版 Pydantic v2，可將 .dict() 改為 .model_dump()，若在 Render 報錯請留意此處)
    df = pd.DataFrame([item.dict() for item in data]).fillna('')

    # 2. 建立 PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    G = nx.Graph()
    edges_info = defaultdict(list)

    # 3. 解析資料與建立關聯
    for _, row in df.iterrows():
        site_a, site_b = str(row['SITE_ID_A']), str(row['SITE_ID_B'])
        p_type, c_id = str(row['PRODUCT_TYPE']).upper(), str(row['CIRCUIT_ID'])
        role_a = 'datacenter' if any(k in str(row['BUILDING_A']) for k in ['HQ', 'DC', '總部', '機房']) else 'customer'
        G.add_node(site_a, label=row['BUILDING_A'], role=role_a)
        circuit_label = f"{p_type} [{c_id}]"
        
        is_cloud_product = any(x in p_type for x in ['MPLS', 'VPN', 'ADSL'])
        if is_cloud_product:
            cloud_id = "Cloud_Network"
            G.add_node(cloud_id, label="MPLS / VPN / Internet", role='cloud')
            edges_info[tuple(sorted((site_a, cloud_id)))].append(circuit_label)
            if site_b:
                role_b = 'datacenter' if any(k in str(row['BUILDING_B']) for k in ['HQ', 'DC', '總部', '機房']) else 'customer'
                G.add_node(site_b, label=row['BUILDING_B'], role=role_b)
                edges_info[tuple(sorted((site_b, cloud_id)))].append(circuit_label)
        elif site_b: 
            role_b = 'datacenter' if any(k in str(row['BUILDING_B']) for k in ['HQ', 'DC', '總部', '機房']) else 'customer'
            G.add_node(site_b, label=row['BUILDING_B'], role=role_b)
            edges_info[tuple(sorted((site_a, site_b)))].append(circuit_label)

    # 4. 取得排版座標
    layout_pos = apply_l2r_topology_layout(G)
    
    # 5. 繪製連線與「線路方框」
    for (u, v), circuits in edges_info.items():
        if u not in layout_pos or v not in layout_pos: continue
        
        x1, y1 = layout_pos[u]
        x2, y2 = layout_pos[v]
        
        # 畫直線
        conn = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.width = Pt(1.5)
        conn.line.color.rgb = RGBColor(117, 117, 117)
        
        mid_x = (x1 + x2) / 2
        mid_y = (y1 + y2) / 2
        box_w_in = 1.6
        box_h_in = 0.25 + 0.15 * len(circuits)
        
        circuit_text = "\n".join(circuits)
        label_box = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mid_x - box_w_in / 2), Inches(mid_y - box_h_in / 2), Inches(box_w_in), Inches(box_h_in))
        label_box.fill.solid()
        label_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        label_box.line.color.rgb = RGBColor(158, 158, 158)
        label_box.text_frame.text = circuit_text
        for p in label_box.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(33, 33, 33)
            p.alignment = PP_ALIGN.CENTER

    # 6. 繪製實體節點
    for node, (x, y) in layout_pos.items():
        node_data = G.nodes[node]
        role = node_data.get('role')
        label = node_data.get('label', node)

        if role == 'cloud':
            w_in, h_in = 2.5, 1.5
            shape = shapes.add_shape(MSO_SHAPE.CLOUD, Inches(x - w_in/2), Inches(y - h_in/2), Inches(w_in), Inches(h_in))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(227, 242, 253)
            shape.line.color.rgb = RGBColor(33, 150, 243)
        elif role == 'datacenter':
            w_in, h_in = 1.8, 0.8
            shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - w_in/2), Inches(y - h_in/2), Inches(w_in), Inches(h_in))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 224, 178)
            shape.line.color.rgb = RGBColor(245, 124, 0)
        else:
            w_in, h_in = 1.8, 0.8
            shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - w_in/2), Inches(y - h_in/2), Inches(w_in), Inches(h_in))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(238, 238, 238)
            shape.line.color.rgb = RGBColor(117, 117, 117)

        shape.text_frame.text = f"{label}\n[{node}]"
        for p in shape.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.CENTER

    # 7. 將生成的 PPT 存入記憶體中
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    # 8. 回傳檔案給 Power Automate
    return Response(
        content=ppt_stream.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": "attachment; filename=network_diagram.pptx"}
    )